import requests
from django.conf import settings
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

FULL_TEXT_PAGE_LIMIT = 30
WORDS_PER_PAGE = 500
FALLBACK_MAX_CHARS = 8000


def extract_text(file_field):
    name = file_field.name.lower()
    file_field.open("rb")
    try:
        if name.endswith(".pdf"):
            reader = PdfReader(file_field)
            page_count = len(reader.pages)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif name.endswith(".docx"):
            document = Document(file_field)
            text = "\n".join(p.text for p in document.paragraphs)
            page_count = max(1, len(text.split()) // WORDS_PER_PAGE)
        elif name.endswith(".pptx"):
            presentation = Presentation(file_field)
            page_count = len(presentation.slides)
            chunks = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
            text = "\n".join(chunks)
        else:
            text = file_field.read().decode("utf-8", errors="ignore")
            page_count = max(1, len(text.split()) // WORDS_PER_PAGE)
    finally:
        file_field.close()

    if page_count <= FULL_TEXT_PAGE_LIMIT:
        return text
    return text[:FALLBACK_MAX_CHARS]


def _call_groq(prompt, system):
    response = requests.post(
        url="https://api.groq.com/openai/v1/chat/completions",
        headers={"Authorization": f"Bearer {settings.GROQ_API_KEY}"},
        json={
            "model": settings.GROQ_MODEL,
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": prompt},
            ],
        },
        timeout=60,
    )
    response.raise_for_status()
    return response.json()["choices"][0]["message"]["content"]


def _call_ollama(prompt, system):
    response = requests.post(
        url=f"{settings.OLLAMA_BASE_URL}/api/chat",
        json={
            "model": settings.OLLAMA_MODEL,
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": prompt},
            ],
            "stream": False,
        },
        timeout=180,
    )
    response.raise_for_status()
    return response.json()["message"]["content"]


def _call_openrouter(prompt, system):
    response = requests.post(
        url="https://openrouter.ai/api/v1/chat/completions",
        headers={"Authorization": f"Bearer {settings.OPENROUTER_API_KEY}"},
        json={
            "model": settings.OPENROUTER_MODEL,
            "messages": [
                {"role": "system", "content": system},
                {"role": "user", "content": prompt},
            ],
        },
        timeout=30,
    )
    response.raise_for_status()
    return response.json()["choices"][0]["message"]["content"]


def call_ai(prompt, system=""):
    if settings.GROQ_API_KEY:
        try:
            return _call_groq(prompt, system)
        except requests.exceptions.RequestException:
            pass

    try:
        return _call_ollama(prompt, system)
    except requests.exceptions.RequestException:
        pass

    if settings.OPENROUTER_API_KEY:
        return _call_openrouter(prompt, system)

    return f"[AI demo javobi, hech qanday AI provayder ishlamayapti]\n\n{prompt[:500]}"
